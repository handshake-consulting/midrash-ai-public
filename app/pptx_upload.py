""" PPTX utilities """
from pptx import Presentation
from app.openai_utils import num_tokens_from_string


def get_presentation(file_path):
    """ PPTX presentation """
    with open(file_path, "rb") as infile:
        return Presentation(infile)

def create_document_dict_pptx(document_path, document_name):
    """ document dict from pptx, used for upserting a pptx """
    document_dict = {}
    presentation = get_presentation(document_path)
    document_dict['title'] = document_name
    text = ""
    start_index = 0
    for slide_index, slide in enumerate(presentation.slides):
        slide_text = ''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += "\n" + shape.text
        text += slide_text
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text += "\n" + text_frame.text
        if num_tokens_from_string(text, "cl100k_base") >= 6500:
            document_dict[start_index] = {}
            document_dict[start_index]['content'] = text
            start_index = slide_index
            text = ""
    return document_dict
